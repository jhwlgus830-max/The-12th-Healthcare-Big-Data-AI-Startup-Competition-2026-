from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# =========================================================
# PRESENTATION SETUP
# =========================================================

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])

# =========================================================
# THEME COLORS
# =========================================================

BG = RGBColor(250, 248, 245)
LAYER2 = RGBColor(235, 232, 244)
LAYER3 = RGBColor(230, 244, 240)
GENERATOR = RGBColor(235, 245, 236)
NAVY = RGBColor(30, 45, 78)
GOLD = RGBColor(245, 158, 11)
BLACK = RGBColor(30, 30, 30)
WHITE = RGBColor(255, 255, 255)

# =========================================================
# BACKGROUND
# =========================================================

slide.background.fill.solid()
slide.background.fill.fore_color.rgb = BG

# =========================================================
# HELPERS
# =========================================================

FONT_NAME = "Malgun Gothic"


def add_round_box(name, x, y, w, h,
                  text,
                  fill,
                  font_size=12,
                  bold=False,
                  line_color=BLACK):

    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y),
        Inches(w), Inches(h)
    )

    shape.name = name

    shape.fill.solid()
print(f"PPT 생성 완료: {output_file}")